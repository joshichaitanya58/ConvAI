

from flask import Flask, render_template, request, redirect, url_for, flash, jsonify, session
from authlib.integrations.flask_client import OAuth
from flask_mail import Mail, Message
from itsdangerous import URLSafeTimedSerializer
from flask_pymongo import PyMongo
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from bson.objectid import ObjectId
import os
import re
import datetime
from dotenv import load_dotenv 
from groq import Groq
from apscheduler.schedulers.background import BackgroundScheduler
from datetime import timedelta, UTC
from google import genai
from google.genai import types
import PyPDF2

# --- App and DB Setup ---
app = Flask(__name__)
load_dotenv() # Load environment variables from .env file

UPLOAD_FOLDER = 'static/uploads/avatars'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
try:
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
except Exception as e:
    print(f"Warning: Could not create upload directory: {e}")

# Load configuration from environment variables
app.config["SECRET_KEY"] = os.environ.get("SECRET_KEY", "your-default-secret-key-for-dev")
app.config["MONGO_URI"] = os.environ.get("MONGO_URI")
app.config['MAIL_SERVER'] = os.environ.get('MAIL_SERVER')
app.config['MAIL_PORT'] = int(os.environ.get('MAIL_PORT', 587))
app.config['MAIL_USE_TLS'] = os.environ.get('MAIL_USE_TLS', 'true').lower() in ['true', 'on', '1']
app.config['MAIL_USERNAME'] = os.environ.get('MAIL_USERNAME')
app.config['MAIL_PASSWORD'] = os.environ.get('MAIL_PASSWORD')
app.config['MAIL_DEFAULT_SENDER'] = os.environ.get('MAIL_DEFAULT_SENDER')

# --- Google Gen AI Client Setup ---
gemini_client = None
_last_api_key = None

def get_gemini_client():
    global gemini_client, _last_api_key
    # Removed load_dotenv() from here to reduce I/O overhead on every API request.
    # Env variables are already loaded at startup.
    current_key = os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")
    
    if gemini_client is None or current_key != _last_api_key:
        _last_api_key = current_key
        if current_key:
            print("Initializing Gemini Client using API Key (Google AI Studio)...")
            gemini_client = genai.Client(api_key=current_key, vertexai=False)
        else:
            print("Initializing Gemini Client using Vertex AI...")
            gemini_client = genai.Client(
                vertexai=True,
                project=os.environ.get("GOOGLE_CLOUD_PROJECT"),
                location=os.environ.get("GOOGLE_CLOUD_LOCATION"),
            )
    return gemini_client

mongo = PyMongo(app)
mail = Mail(app)
oauth = OAuth(app)

# --- Google OAuth Config ---
google = oauth.register(
    name='google',
    client_id=os.environ.get("GOOGLE_CLIENT_ID"),
    client_secret=os.environ.get("GOOGLE_CLIENT_SECRET"),
    # Use server metadata for automatic discovery of endpoints (including jwks_uri)
    server_metadata_url='https://accounts.google.com/.well-known/openid-configuration',
    client_kwargs={'scope': 'openid email profile'},
    api_base_url='https://www.googleapis.com/oauth2/v3/',
    userinfo_endpoint='https://openidconnect.googleapis.com/v1/userinfo',
)

# --- Serializer for password reset tokens ---
s = URLSafeTimedSerializer(app.config['SECRET_KEY'])

# --- Flask-Login Setup ---
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login' # Redirect to /login if user is not authenticated

# --- User Model ---
class User(UserMixin):
    def __init__(self, user_data):
        self.id = str(user_data["_id"])
        self.email = user_data.get("email") or user_data.get("username")
        self.username = user_data["username"]
        self.password_hash = user_data["password"]
        self.theme = user_data.get("theme", "dark") # Default to dark theme
        self.is_premium = user_data.get("is_premium", False)
        self.font_size = user_data.get("font_size", "medium")
        self.plan_type = user_data.get("plan_type")
        self.subscription_start_date = user_data.get("subscription_start_date")
        self.avatar_url = user_data.get("avatar_url")
        self.is_oauth_user = user_data.get("is_oauth_user", False)

    @staticmethod
    def get(user_id):
        try:
            user_data = mongo.db.users.find_one({"_id": ObjectId(user_id)})
            if user_data:
                return User(user_data)
        except Exception:
            pass
        return None

@login_manager.user_loader
def load_user(user_id):
    return User.get(user_id)

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.template_filter('format_datetime')
def _jinja2_filter_format_datetime(value, format='%B %d, %Y'):
    """A custom Jinja2 filter to format a datetime object or the string 'now'."""
    if value == "now":
        dt = datetime.datetime.now(datetime.UTC)
    elif isinstance(value, datetime.datetime):
        dt = value
    else:
        return value
    return dt.strftime(format)

# --- Routes ---
@app.route('/')
@login_required
def index():
    # current_user is available in templates automatically
    return render_template('index.html')

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        email = request.form.get('email', '').lower().strip()
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')

        if not all([username, email, password, confirm_password]):
            flash('All fields are required.', 'error')
            return render_template('signup.html', form_data=request.form)

        if password != confirm_password:
            flash('Passwords do not match. Please try again.', 'error')
            return render_template('signup.html', form_data=request.form)

        # Password strength check (simple version)
        if len(password) < 8:
            flash('Password must be at least 8 characters long.', 'error')
            return render_template('signup.html', form_data=request.form)

        if mongo.db.users.find_one({'username': {'$regex': f'^{re.escape(username)}$', '$options': 'i'}}):
            flash('Username already exists. Please choose a different one.', 'error')
            return render_template('signup.html', form_data=request.form)

        # Check if email exists in active users
        existing_user = mongo.db.users.find_one({'email': {'$regex': f'^{re.escape(email)}$', '$options': 'i'}})
        if existing_user:
            flash('An account with this email already exists.', 'error')
            return redirect(url_for('signup'))

        # --- Reactivation Logic ---
        # Check if email exists in deleted users
        deleted_user_data = mongo.db.deleted_users.find_one({'email': email})
        is_reactivating = False
        if deleted_user_data:
            is_reactivating = True
            # Use the original username for consistency
            username = deleted_user_data.get('username', username)
            flash_message = 'Welcome back! A verification code has been sent to reactivate your account.'
            email_subject = 'Reactivate Your ConvAI Account'
        else:
            flash_message = 'A 6-digit verification code has been sent to your email.'
            email_subject = 'Verify Your ConvAI Account'

        hashed_password = generate_password_hash(password, method='pbkdf2:sha256')

        # Generate and send OTP
        import random
        otp = f"{random.randint(100000, 999999)}"
        otp_expiry = datetime.datetime.now(UTC) + datetime.timedelta(minutes=10)
        session['signup_verification'] = {
            'username': username,
            'email': email,
            'password': hashed_password,
            'otp': otp,
            'otp_expiry': otp_expiry.isoformat(),
            'is_reactivating': is_reactivating # Add reactivation flag
        }
        send_otp_email(email, otp, email_subject)

        flash(flash_message, 'success')
        return redirect(url_for('verify_otp', flow='signup'))

    return render_template('signup.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password')
        
        # Find user by username or email (case-insensitive)
        escaped_username = re.escape(username)
        user_data = mongo.db.users.find_one({
            '$or': [
                {'username': {'$regex': f'^{escaped_username}$', '$options': 'i'}},
                {'email': {'$regex': f'^{escaped_username}$', '$options': 'i'}}
            ]
        })

        if user_data and check_password_hash(user_data['password'], password):
            user = User(user_data)

            # --- 2FA Logic ---
            import random
            otp = f"{random.randint(100000, 999999)}"
            otp_expiry = datetime.datetime.now(UTC) + datetime.timedelta(minutes=10)

            session['login_verification'] = {
                'user_id': str(user.id),
                'otp': otp,
                'otp_expiry': otp_expiry.isoformat()
            }

            send_otp_email(user.email, otp, "Your ConvAI Login Code")

            flash('A 6-digit login code has been sent to your email.', 'success')
            return redirect(url_for('verify_otp', flow='login'))
            # --- End 2FA Logic ---

        flash('Invalid username or password. Please try again.', 'error')
        return redirect(url_for('login'))

    return render_template('login.html')

@app.route('/google/login')
def google_login():
    """Redirects to Google's authentication page."""
    # The redirect_uri is automatically detected by authlib if not specified.
    # It's often more reliable to let the library handle it.
    # The URI must still be configured in the Google Cloud Console.
    return google.authorize_redirect(url_for('google_auth', _external=True))

@app.route('/google/auth')
def google_auth():
    """Callback route for Google to redirect to."""
    try:
        token = google.authorize_access_token()
        # The 'userinfo' is a standard claim in OpenID Connect,
        # and authlib can parse it from the ID token for you.
        user_info = token.get('userinfo')
        email = user_info.get('email') if user_info else None

        if not email:
            flash('Could not retrieve email from Google. Please try again.', 'error')
            return redirect(url_for('login'))

        # Find or create user
        user_data = mongo.db.users.find_one({'email': email})

        if user_data:
            # User exists, log them in
            user = User(user_data)
            login_user(user)
            return redirect(url_for('index'))
        else:
            # User does not exist, create a new account
            # We use a placeholder for the password as it won't be used for login
            hashed_password = generate_password_hash(os.urandom(24).hex(), method='pbkdf2:sha256')
            new_user_data = {
                'username': user_info.get('name', email.split('@')[0]),
                'email': email,
                'password': hashed_password,
                # Save the Google profile picture URL
                'avatar_url': user_info.get('picture'), 
                'is_oauth_user': True # Flag to identify Google users
            }
            mongo.db.users.insert_one(new_user_data)
            user = User(new_user_data)
            login_user(user)
            return redirect(url_for('index'))
    except Exception as e:
        flash(f'An error occurred during Google authentication: {e}', 'error')
        return redirect(url_for('login'))

@app.route('/verify-otp/<flow>', methods=['GET', 'POST'])
def verify_otp(flow):
    """Handles OTP verification for both signup and login flows."""
    if flow not in ['signup', 'login']:
        return redirect(url_for('login'))

    session_key = f'{flow}_verification'
    verification_data = session.get(session_key)

    if not verification_data:
        flash('Verification session expired. Please try again.', 'error')
        return redirect(url_for('signup' if flow == 'signup' else 'login'))

    email = verification_data.get('email')
    if not email and flow == 'login':
        user = User.get(verification_data.get('user_id'))
        if user:
            email = user.email

    if request.method == 'POST':
        submitted_otp = request.form.get('otp', '').strip()
        
        # Validate OTP
        if submitted_otp != verification_data.get('otp'):
            flash('Invalid OTP. Please try again.', 'error')
            return render_template('verify_otp.html', flow=flow, email=email)

        # Validate Expiry
        expiry_time = datetime.datetime.fromisoformat(verification_data.get('otp_expiry'))
        if datetime.datetime.now(UTC).replace(tzinfo=None) > expiry_time.replace(tzinfo=None):
            flash('OTP has expired. Please request a new one.', 'error')
            return render_template('verify_otp.html', flow=flow, email=email)

        # --- Process based on flow ---
        if flow == 'signup':
            # Create user account
            is_reactivating = verification_data.get('is_reactivating', False)
            username = verification_data['username']
            email = verification_data['email']

            mongo.db.users.insert_one({
                'username': username,
                'email': email,
                'password': verification_data['password']
            })

            if is_reactivating:
                # --- Handle Reactivation ---
                # 1. Delete from the deleted_users collection
                mongo.db.deleted_users.delete_one({'email': email})
                # 2. Send "Welcome Back" email
                try:
                    msg = Message('Welcome Back to ConvAI!', recipients=[email])
                    msg.html = render_template('email_welcome_back.html', username=username)
                    mail.send(msg)
                    print(f"Sent reactivation (welcome back) email to {email}")
                except Exception as e:
                    print(f"Failed to send reactivation email to {email}: {e}")
            else:
                # --- Handle New User ---
                # Send standard welcome email
                try:
                    msg = Message('Welcome to ConvAI!', recipients=[email])
                    msg.html = render_template('email_welcome.html', username=username)
                    mail.send(msg)
                except Exception as e:
                    print(f"Failed to send welcome email to {email}: {e}")

            session.pop(session_key, None)
            flash('Your account has been successfully created! Please log in.', 'success')
            return redirect(url_for('login'))

        elif flow == 'login':
            user = User.get(verification_data['user_id'])
            if user:
                login_user(user)
                session.pop(session_key, None)
                return redirect(url_for('index'))
            else:
                flash('Could not find user to log in. Please try again.', 'error')
                return redirect(url_for('login'))

    return render_template('verify_otp.html', flow=flow, email=email)

@app.route('/resend-otp/<flow>', methods=['POST'])
def resend_verification_otp(flow):
    if flow not in ['signup', 'login']:
        return jsonify({"success": False, "error": "Invalid flow"}), 400

    session_key = f'{flow}_verification'
    verification_data = session.get(session_key)

    if not verification_data:
        return jsonify({"success": False, "error": "Verification session not found or expired."}), 404

    email = verification_data.get('email')
    if not email and flow == 'login':
        user = User.get(verification_data.get('user_id'))
        if user:
            email = user.email

    if not email:
        return jsonify({"success": False, "error": "Could not determine email address."}), 400

    # Generate and send new OTP
    import random
    otp = f"{random.randint(100000, 999999)}"
    otp_expiry = datetime.datetime.now(UTC) + datetime.timedelta(minutes=10)
    
    # Update session with new OTP
    verification_data['otp'] = otp
    verification_data['otp_expiry'] = otp_expiry.isoformat()
    session[session_key] = verification_data

    subject = "Your New ConvAI Verification Code" if flow == 'signup' else "Your New ConvAI Login Code"
    send_otp_email(email, otp, subject)

    return jsonify({"success": True})

def send_otp_email(email, otp, subject):
    """Helper function to send OTP emails."""
    try:
        msg = Message(subject, recipients=[email])
        msg.html = render_template('email_otp.html', otp=otp, subject=subject)
        mail.send(msg)
        print(f"OTP email sent to {email} with subject: '{subject}'")
    except Exception as e:
        print(f"Failed to send OTP email to {email}: {e}")


@app.route('/forgot-password', methods=['GET', 'POST'])
def forgot_password():
    if request.method == 'POST':
        email = request.form.get('email', '').lower().strip()
        user_data = mongo.db.users.find_one({"email": email})

        if user_data:
            import random
            otp = f"{random.randint(100000, 999999)}"
            otp_expiry = datetime.datetime.now(UTC) + datetime.timedelta(minutes=10)

            mongo.db.users.update_one(
                {"_id": user_data["_id"]},
                {"$set": {
                    "reset_otp": otp,
                    "reset_otp_expiry": otp_expiry
                }}
            )

            reset_url = url_for('reset_password', email=email, _external=True)

            try:
                msg = Message('Password Reset OTP for ConvAI', recipients=[email])
                msg.body = (
                    f"Hello,\n\n"
                    f"We received a request to reset the password for your ConvAI account.\n"
                    f"Your One-Time Password (OTP) is:\n\n"
                    f"{otp}\n\n"
                    f"This code is valid for 10 minutes. If you did not request this, you can safely ignore this email.\n\n"
                    f"Reset Page: {reset_url}\n\n"
                    f"- ConvAI Team"
                )
                msg.html = f"""
<html><body style="margin:0;padding:0;background-color:#0e1116;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,Helvetica,Arial,sans-serif;">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="background-color:#0e1116;padding:40px 20px;"><tr><td align="center">
<table role="presentation" width="600" cellpadding="0" cellspacing="0" style="max-width:600px;width:100%;">
  <tr><td align="center" style="padding-bottom:32px;">
    <table role="presentation" cellpadding="0" cellspacing="0"><tr>
      <td style="width:36px;height:36px;border-radius:9px;background:linear-gradient(150deg,#f0bd76,#e0a458 70%);text-align:center;vertical-align:middle;font-family:'Georgia',serif;font-weight:bold;font-size:18px;color:#0e1116;">C</td>
      <td style="padding-left:10px;font-family:'Georgia',serif;font-size:20px;font-weight:bold;color:#f3efe6;letter-spacing:0.3px;">ConvAI</td>
    </tr></table>
  </td></tr>
  <tr><td style="background-color:#14181f;border:1px solid #1f252e;border-radius:16px;overflow:hidden;">
    <table role="presentation" width="100%" cellpadding="0" cellspacing="0">
      <tr><td style="height:4px;background:linear-gradient(90deg,#e0a458,#f0bd76,#e0a458);"></td></tr>
    </table>
    <table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="padding:36px 32px;">
      <tr><td style="font-family:'Georgia',serif;font-size:22px;font-weight:bold;color:#f3efe6;text-align:center;padding-bottom:8px;">Password Reset Request</td></tr>
      <tr><td style="font-size:15px;color:#cfc9bc;text-align:center;padding-bottom:28px;line-height:1.6;">We received a request to reset the password for your ConvAI account. Use the code below to proceed.</td></tr>
      <tr><td align="center" style="padding-bottom:24px;">
        <div style="background-color:#191e26;border:1px solid #1f252e;border-radius:12px;padding:20px 32px;display:inline-block;">
          <span style="font-family:'Courier New',monospace;font-size:32px;font-weight:bold;letter-spacing:8px;color:#e0a458;">{otp}</span>
        </div>
      </td></tr>
      <tr><td style="font-size:13px;color:#6b7585;text-align:center;padding-bottom:20px;line-height:1.6;">This code is valid for <strong style="color:#cfc9bc;">10 minutes</strong>. If you did not request this, you can safely ignore this email.</td></tr>
      <tr><td align="center" style="padding-bottom:8px;"><a href="{reset_url}" style="display:inline-block;background:#e0a458;color:#0e1116;text-decoration:none;padding:12px 32px;border-radius:10px;font-weight:bold;font-size:14px;">Reset My Password →</a></td></tr>
    </table>
  </td></tr>
  <tr><td style="padding-top:28px;text-align:center;">
    <p style="font-size:12px;color:#6b7585;margin:0 0 4px;line-height:1.5;">This is an automated message from ConvAI. Please do not reply.</p>
    <p style="font-size:12px;color:#6b7585;margin:0;line-height:1.5;">Need help? Contact us at <a href="mailto:joshichaitanya58@gmail.com" style="color:#e0a458;text-decoration:none;">joshichaitanya58@gmail.com</a></p>
  </td></tr>
</table>
</td></tr></table>
</body></html>
"""
                mail.send(msg)
                print(f"Password reset OTP sent to {email}. Reset URL: {reset_url}")
                flash('A 6-digit OTP code has been sent to your email.', 'success')
                session['reset_email'] = email
                return redirect(url_for('reset_password'))
            except Exception as e:
                flash('Failed to send email. Please check your configuration.', 'error')
                print(f"Email sending error: {e}")
        else:
            flash('No account found with that email address.', 'error')

        return redirect(url_for('forgot_password'))

    return render_template('forgot_password.html')

@app.route('/reset-password', methods=['GET', 'POST'])
def reset_password():
    email = request.args.get('email', '').strip()
    if not email:
        email = session.get('reset_email', '')

    if request.method == 'POST':
        email = request.form.get('email', '').lower().strip()
        otp = request.form.get('otp', '').strip()
        password = request.form.get('password')
        confirm_password = request.form.get('confirm_password')

        if not email:
            flash('Session expired or email missing. Please try again.', 'error')
            return redirect(url_for('forgot_password'))

        if password != confirm_password:
            flash('Passwords do not match.', 'error')
            return redirect(url_for('reset_password', email=email))

        user_data = mongo.db.users.find_one({'email': email})
        if not user_data:
            flash('User not found.', 'error')
            return redirect(url_for('forgot_password'))

        db_otp = user_data.get('reset_otp')
        db_expiry = user_data.get('reset_otp_expiry')

        if not db_otp or not db_expiry:
            flash('No password reset requested or OTP expired.', 'error')
            return redirect(url_for('forgot_password'))

        now = datetime.datetime.now(UTC)
        if db_expiry.tzinfo is None:
            now = now.replace(tzinfo=None)

        if db_otp != otp:
            flash('Invalid OTP code. Please try again.', 'error')
            return redirect(url_for('reset_password', email=email))

        if now > db_expiry:
            flash('The OTP code has expired. Please request a new one.', 'error')
            return redirect(url_for('forgot_password'))

        hashed_password = generate_password_hash(password, method='pbkdf2:sha256')
        mongo.db.users.update_one(
            {'_id': user_data['_id']},
            {
                '$set': {'password': hashed_password},
                '$unset': {'reset_otp': '', 'reset_otp_expiry': ''}
            }
        )

        session.pop('reset_email', None)
        flash('Your password has been updated successfully! You can now log in.', 'success')
        return redirect(url_for('login'))

    if not email:
        flash('Please request a password reset first.', 'error')
        return redirect(url_for('forgot_password'))

    return render_template('reset_password.html', email=email)

@app.route('/api/resend-otp', methods=['POST'])
def resend_otp():
    data = request.get_json() or {}
    email = data.get('email', '').lower().strip()
    if not email:
        email = session.get('reset_email', '')

    if not email:
        return jsonify({"success": False, "error": "Email is required"}), 400

    user_data = mongo.db.users.find_one({"email": email})
    if not user_data:
        return jsonify({"success": False, "error": "No account found with that email address"}), 404

    import random
    otp = f"{random.randint(100000, 999999)}"
    otp_expiry = datetime.datetime.now(UTC) + datetime.timedelta(minutes=10)

    mongo.db.users.update_one(
        {"_id": user_data["_id"]},
        {"$set": {
            "reset_otp": otp,
            "reset_otp_expiry": otp_expiry
        }}
    )

    reset_url = url_for('reset_password', email=email, _external=True)

    try:
        msg = Message('Password Reset OTP for ConvAI (Resend)', recipients=[email])
        msg.body = (
            f"Hello,\n\n"
            f"We received a request to resend the password reset OTP for your ConvAI account.\n"
            f"Your new One-Time Password (OTP) is:\n\n"
            f"{otp}\n\n"
            f"This code is valid for 10 minutes. If you did not request this, you can safely ignore this email.\n\n"
            f"Reset Page: {reset_url}\n\n"
            f"- ConvAI Team"
        )
        msg.html = f"""
<html><body style="margin:0;padding:0;background-color:#0e1116;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Roboto,Helvetica,Arial,sans-serif;">
<table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="background-color:#0e1116;padding:40px 20px;"><tr><td align="center">
<table role="presentation" width="600" cellpadding="0" cellspacing="0" style="max-width:600px;width:100%;">
  <tr><td align="center" style="padding-bottom:32px;">
    <table role="presentation" cellpadding="0" cellspacing="0"><tr>
      <td style="width:36px;height:36px;border-radius:9px;background:linear-gradient(150deg,#f0bd76,#e0a458 70%);text-align:center;vertical-align:middle;font-family:'Georgia',serif;font-weight:bold;font-size:18px;color:#0e1116;">C</td>
      <td style="padding-left:10px;font-family:'Georgia',serif;font-size:20px;font-weight:bold;color:#f3efe6;letter-spacing:0.3px;">ConvAI</td>
    </tr></table>
  </td></tr>
  <tr><td style="background-color:#14181f;border:1px solid #1f252e;border-radius:16px;overflow:hidden;">
    <table role="presentation" width="100%" cellpadding="0" cellspacing="0">
      <tr><td style="height:4px;background:linear-gradient(90deg,#e0a458,#f0bd76,#e0a458);"></td></tr>
    </table>
    <table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="padding:36px 32px;">
      <tr><td style="font-family:'Georgia',serif;font-size:22px;font-weight:bold;color:#f3efe6;text-align:center;padding-bottom:8px;">New OTP Code</td></tr>
      <tr><td style="font-size:15px;color:#cfc9bc;text-align:center;padding-bottom:28px;line-height:1.6;">Here is your new One-Time Password to reset your ConvAI account password.</td></tr>
      <tr><td align="center" style="padding-bottom:24px;">
        <div style="background-color:#191e26;border:1px solid #1f252e;border-radius:12px;padding:20px 32px;display:inline-block;">
          <span style="font-family:'Courier New',monospace;font-size:32px;font-weight:bold;letter-spacing:8px;color:#e0a458;">{otp}</span>
        </div>
      </td></tr>
      <tr><td style="font-size:13px;color:#6b7585;text-align:center;padding-bottom:20px;line-height:1.6;">This code is valid for <strong style="color:#cfc9bc;">10 minutes</strong>. If you did not request this, you can safely ignore this email.</td></tr>
      <tr><td align="center" style="padding-bottom:8px;"><a href="{reset_url}" style="display:inline-block;background:#e0a458;color:#0e1116;text-decoration:none;padding:12px 32px;border-radius:10px;font-weight:bold;font-size:14px;">Reset My Password →</a></td></tr>
    </table>
  </td></tr>
  <tr><td style="padding-top:28px;text-align:center;">
    <p style="font-size:12px;color:#6b7585;margin:0 0 4px;line-height:1.5;">This is an automated message from ConvAI. Please do not reply.</p>
    <p style="font-size:12px;color:#6b7585;margin:0;line-height:1.5;">Need help? Contact us at <a href="mailto:joshichaitanya58@gmail.com" style="color:#e0a458;text-decoration:none;">joshichaitanya58@gmail.com</a></p>
  </td></tr>
</table>
</td></tr></table>
</body></html>
"""
        mail.send(msg)
        print(f"Resent OTP to {email}: {otp}")
        return jsonify({"success": True})
    except Exception as e:
        print(f"Email resending error: {e}")
        return jsonify({"success": False, "error": "Failed to send email. Please check your configuration."}), 500


@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

@app.route('/profile', methods=['GET', 'POST'])
@login_required
def profile():
    if request.method == 'POST':
        # Check which form was submitted
        if 'update_details' in request.form:
            username = request.form.get('username')
            email = request.form.get('email', '').lower().strip()

            # Basic validation
            if not username or not email:
                flash('Username and email are required.', 'error')
                return redirect(url_for('profile'))

            # Check if new username is already taken by another user
            if username != current_user.username and mongo.db.users.find_one({'username': {'$regex': f'^{re.escape(username)}$', '$options': 'i'}}):
                flash('Username already exists. Please choose a different one.', 'error')
                return redirect(url_for('profile'))

            # Check if new email is already taken by another user
            if email != current_user.email and mongo.db.users.find_one({'email': {'$regex': f'^{re.escape(email)}$', '$options': 'i'}}):
                flash('An account with this email already exists.', 'error')
                return redirect(url_for('profile'))

            update_fields = {
                'username': username,
                'email': email
            }

            # Handle avatar upload
            if 'avatar' in request.files:
                file = request.files['avatar']
                if file and file.filename != '' and allowed_file(file.filename):
                    # Create a secure, unique filename
                    filename = secure_filename(file.filename)
                    file_ext = filename.rsplit('.', 1)[1].lower()
                    unique_filename = f"{current_user.id}.{file_ext}"
                    
                    # Delete old avatar if it exists and has a different extension
                    if current_user.avatar_url:
                        old_path = os.path.join(app.root_path, current_user.avatar_url.lstrip('/'))
                        if os.path.exists(old_path) and os.path.basename(old_path) != unique_filename:
                            os.remove(old_path)

                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                    file.save(file_path)
                    update_fields['avatar_url'] = f"/{file_path.replace(os.sep, '/')}"

            mongo.db.users.update_one(
                {"_id": ObjectId(current_user.id)},
                {"$set": update_fields}
            )
            flash('Profile details updated successfully!', 'success')

        elif 'delete_avatar' in request.form:
            if current_user.avatar_url:
                # Delete physical file if local
                if not current_user.avatar_url.startswith('http'):
                    try:
                        file_path = os.path.join(app.root_path, current_user.avatar_url.lstrip('/'))
                        if os.path.exists(file_path):
                            os.remove(file_path)
                    except Exception as e:
                        print(f"Error removing avatar file: {e}")
                
                # Update DB to unset the avatar_url
                mongo.db.users.update_one(
                    {"_id": ObjectId(current_user.id)},
                    {"$unset": {"avatar_url": ""}}
                )
                flash('Profile picture removed successfully!', 'success')
            return redirect(url_for('profile'))

        elif 'update_appearance' in request.form:
            if not current_user.is_premium:
                flash('This feature is for premium users only.', 'error')
                return redirect(url_for('profile'))
            
            font_size = request.form.get('font_size')
            if font_size in ['small', 'medium', 'large']:
                mongo.db.users.update_one({"_id": ObjectId(current_user.id)}, {"$set": {"font_size": font_size}})
                flash('Appearance settings updated!', 'success')

        elif 'change_password' in request.form:
            current_password = request.form.get('current_password')
            new_password = request.form.get('new_password')
            confirm_password = request.form.get('confirm_password')

            if not all([current_password, new_password, confirm_password]):
                flash('All password fields are required.', 'error')
                return redirect(url_for('profile'))

            if not check_password_hash(current_user.password_hash, current_password):
                flash('Your current password is not correct.', 'error')
                return redirect(url_for('profile'))

            if new_password != confirm_password:
                flash('New passwords do not match.', 'error')
                return redirect(url_for('profile'))

            hashed_password = generate_password_hash(new_password, method='pbkdf2:sha256')
            mongo.db.users.update_one(
                {"_id": ObjectId(current_user.id)},
                {"$set": {'password': hashed_password}}
            )
            flash('Your password has been updated successfully!', 'success')

        return redirect(url_for('profile'))

    return render_template('profile.html')

@app.route('/help')
@login_required
def help_page():
    """Displays the help and support page."""
    return render_template('help.html')

@app.route('/terms')
def terms_page():
    """Displays the terms and conditions page."""
    return render_template('terms.html')

@app.route('/privacy')
def privacy_page():
    """Displays the privacy policy page."""
    return render_template('privacy.html')

@app.route('/plans')
@login_required
def plans_page():
    """Displays the available subscription plans."""
    return render_template('plans.html')

@app.route('/payment/<plan_name>')
@login_required
def payment_page(plan_name):
    """Displays a mock payment form for the selected plan with dynamic UPI QR code."""
    if plan_name not in ['monthly', 'yearly']:
        flash('Invalid plan selected.', 'error')
        return redirect(url_for('plans_page'))
    
    upi_id = os.environ.get("UPI_ID", "vaccinetrack17@okaxis")
    monthly_price = os.environ.get("MONTHLY_PRICE", "299")
    yearly_price = os.environ.get("YEARLY_PRICE", "3500")
    
    plan_details = {
        'monthly': {'name': 'Monthly Plan', 'price': f'₹{monthly_price}', 'amount': monthly_price},
        'yearly': {'name': 'Yearly Plan', 'price': f'₹{yearly_price}', 'amount': yearly_price}
    }
    plan_info = plan_details[plan_name]
    
    # Generate UPI URI: upi://pay?pa=UPI_ID&pn=ConvAI&am=Amount&cu=INR
    import urllib.parse
    upi_url = f"upi://pay?pa={upi_id}&pn=ConvAI&am={plan_info['amount']}&cu=INR"
    encoded_upi_url = urllib.parse.quote(upi_url)
    
    # Free QR code generator API
    qr_code_url = f"https://api.qrserver.com/v1/create-qr-code/?size=150x150&data={encoded_upi_url}"
    
    return render_template('payment.html', plan=plan_info, plan_key=plan_name, qr_code_url=qr_code_url, upi_id=upi_id)

@app.route('/payment/success', methods=['POST'])
@login_required
def payment_success():
    """Handles the successful payment and upgrades the user to premium, logging the transaction."""
    plan_type = request.form.get('plan_type') # 'monthly' or 'yearly'
    if not plan_type:
        flash('An error occurred. Invalid plan type.', 'error')
        return redirect(url_for('plans_page'))

    payment_method = request.form.get('payment_method', 'card')
    
    # Securely retrieve details based on payment method
    details = ""
    if payment_method == 'card':
        card_num = request.form.get('card-number', '').strip()
        card_name = request.form.get('card-name', '').strip()
        # Mask the card number to only show last 4 digits
        masked_num = f"xxxx xxxx xxxx {card_num[-4:]}" if len(card_num) >= 4 else "xxxx xxxx xxxx"
        details = f"Card ({masked_num}) - {card_name}"
    elif payment_method == 'upi':
        upi_id = request.form.get('upi-id', '').strip()
        details = f"UPI ID: {upi_id}"
    elif payment_method == 'netbanking':
        bank = request.form.get('bank', '').strip()
        details = f"Netbanking - {bank}"

    # Determine amount based on env variables
    monthly_price = os.environ.get("MONTHLY_PRICE", "299")
    yearly_price = os.environ.get("YEARLY_PRICE", "3500")
    amount = f"₹{monthly_price}" if plan_type == 'monthly' else f"₹{yearly_price}"

    # In a real app, you would verify the payment here with a payment provider (e.g., Stripe)
    now = datetime.datetime.now(UTC)
    mongo.db.users.update_one(
        {"_id": ObjectId(current_user.id)},
        {"$set": {"is_premium": True, "plan_type": plan_type, "subscription_start_date": now}}
    )

    # Log payment transaction in database
    try:
        mongo.db.payments.insert_one({
            "user_id": ObjectId(current_user.id),
            "username": current_user.username,
            "plan_type": plan_type,
            "payment_method": payment_method,
            "payment_details": details,
            "amount": amount,
            "timestamp": now,
            "status": "success"
        })
        print(f"Logged payment transaction for user {current_user.username} ({plan_type}) using {payment_method}")
    except Exception as e:
        print(f"Error logging payment transaction: {e}")

    # Send a confirmation email
    try:
        msg = Message('Your ConvAI Premium Subscription is Active!', recipients=[current_user.email])
        msg.html = render_template('payment_confirmation.html', username=current_user.username)
        mail.send(msg)
    except Exception as e:
        print(f"Error sending payment confirmation email for user {current_user.id}: {e}")
        # Don't block the user flow if email fails, just log it.

    flash('Congratulations! You are now a Premium user.', 'success')
    return redirect(url_for('profile'))

@app.route('/api/user/theme', methods=['POST'])
@login_required
def set_user_theme():
    data = request.get_json()
    theme = data.get('theme')

    if theme not in ['light', 'dark']:
        return jsonify({"success": False, "error": "Invalid theme"}), 400

    mongo.db.users.update_one(
        {"_id": ObjectId(current_user.id)},
        {"$set": {"theme": theme}}
    )

    return jsonify({"success": True})

@app.route('/api/user/appearance', methods=['POST'])
@login_required
def set_user_appearance():
    if not current_user.is_premium:
        return jsonify({"success": False, "error": "Premium only"}), 403

    data = request.get_json()
    font_size = data.get('font_size')

    if font_size not in ['small', 'medium', 'large']:
        return jsonify({"success": False, "error": "Invalid font size"}), 400

    mongo.db.users.update_one(
        {"_id": ObjectId(current_user.id)},
        {"$set": {"font_size": font_size}}
    )
    return jsonify({"success": True})

@app.route('/delete-account', methods=['POST'])
@login_required
def delete_account():
    password = request.form.get('password')

    # Get user document to check if OAuth user
    user_data = mongo.db.users.find_one({"_id": ObjectId(current_user.id)})
    is_oauth_user = user_data.get('is_oauth_user', False) if user_data else False

    # 1. Verify password if not an OAuth user
    if not is_oauth_user:
        if not password or not check_password_hash(current_user.password_hash, password):
            flash('Incorrect password. Account deletion failed.', 'error')
            return redirect(url_for('profile'))

    user_id = ObjectId(current_user.id)

    # Store user details for email before deletion
    user_email_for_notification = user_data.get('email')
    user_name_for_notification = user_data.get('username')
    # New Step: Log user info to a separate collection before deletion
    if user_data:
        mongo.db.deleted_users.insert_one({
            'username': user_data.get('username'),
            'email': user_data.get('email'),
            'deleted_at': datetime.datetime.now(UTC)
        })
        print(f"Archived user info for {user_data.get('email')}")

    # 2. Delete associated data
    # Delete chats
    mongo.db.chats.delete_many({"user_id": user_id})
    # Delete ratings
    mongo.db.ratings.delete_many({"user_id": user_id})

    # 3. Delete avatar file if it exists
    if current_user.avatar_url:
        try:
            avatar_path = os.path.join(app.root_path, current_user.avatar_url.lstrip('/'))
            if os.path.exists(avatar_path):
                os.remove(avatar_path)
        except Exception as e:
            print(f"Error deleting avatar for user {current_user.id}: {e}")

    # 4. Delete the user document
    mongo.db.users.delete_one({"_id": user_id})

    # 5. Log the user out
    logout_user()

    # Send account deletion confirmation email
    if user_email_for_notification:
        try:
            msg = Message('Your ConvAI Account Has Been Deleted', recipients=[user_email_for_notification])
            msg.html = render_template('email_account_deleted.html', username=user_name_for_notification)
            mail.send(msg)
        except Exception as e:
            print(f"Failed to send account deletion email to {user_email_for_notification}: {e}")

    flash('Your account has been permanently deleted.', 'success')
    return redirect(url_for('login'))

@app.route('/.well-known/appspecific/com.chrome.devtools.json')
def chrome_devtools_config():
    """
    Handles a request from Chrome DevTools to prevent 404 errors in logs.
    This is for debugging purposes and is not essential for the app to run.
    """
    return jsonify({})


def resolve_model_name(requested_model):
    """Map friendly UI model names to the actual model ID expected by Groq."""
    requested_model = (requested_model or '').strip()
    if not requested_model:
        return os.environ.get("GROQ_MODEL", "llama-3.3-70b-versatile")

    aliases = {
    "groq/compound-mini": "groq/compound-mini",
    "groq/compound": "groq/compound",
}

    return aliases.get(requested_model, requested_model)


# --- API Routes ---
@app.route('/api/chat', methods=['POST'])
@login_required
def api_chat():
    # Support both JSON (used by regenerate) and FormData (used by normal submit with attachment)
    if request.is_json:
        data = request.get_json() or {}
        user_message = data.get("message")
        chat_id = data.get("chat_id")
        requested_model = data.get("model", "groq/compound-mini")
        file = None
    else:
        user_message = request.form.get("message")
        chat_id = request.form.get("chat_id") # Can be None for a new chat
        requested_model = request.form.get("model", "groq/compound-mini")
        file = request.files.get('file')

    model = resolve_model_name(requested_model)

    if not user_message and not file:
        return jsonify({"error": "No message or file provided"}), 400

    print("-" * 50)
    print(f"Received new chat request from user: {current_user.username}")
    print(f"Requested model (from UI): '{requested_model}'")
    
    # --- File Processing ---
    file_content = ""
    file_prompt_prefix = ""
    if file:
        filename = file.filename or "uploaded_file"
        file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        print(f"Processing uploaded file: {filename} (type: {file_ext})")

        text_based_ext = ['pdf', 'txt', 'md', 'csv', 'docx', 'xlsx', 'xls', 'pptx', 'ppt']
        media_based_ext = ['png', 'jpg', 'jpeg', 'gif', 'mp3', 'mp4', 'wav']

        if file_ext in text_based_ext:
            # For text-based files, extract content
            if file_ext == 'pdf':
                try:
                    pdf_reader = PyPDF2.PdfReader(file.stream)
                    for page in pdf_reader.pages:
                        file_content += page.extract_text()
                except Exception as e:
                    print(f"Error reading PDF: {e}")
                    return jsonify({"reply": f"Sorry, I couldn't read the PDF file '{filename}'. It might be corrupted or protected."}), 400
            elif file_ext in ['txt', 'md', 'csv']:
                file_content = file.read().decode('utf-8', errors='ignore')
            elif file_ext == 'docx':
                try:
                    import docx
                    doc = docx.Document(file.stream)
                    file_content = "\n".join([para.text for para in doc.paragraphs])
                except Exception as e:
                    print(f"Error reading DOCX: {e}")
                    return jsonify({"reply": f"Sorry, I couldn't read the DOCX file '{filename}'."}), 400
            elif file_ext in ['xlsx', 'xls']:
                try:
                    import openpyxl
                    workbook = openpyxl.load_workbook(file.stream)
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        file_content += f"--- Sheet: {sheet_name} ---\n"
                        for row in sheet.iter_rows(values_only=True):
                            file_content += ", ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                except Exception as e:
                    print(f"Error reading Excel file: {e}")
                    return jsonify({"reply": f"Sorry, I couldn't read the Excel file '{filename}'."}), 400
            elif file_ext in ['pptx', 'ppt']:
                try:
                    import pptx
                    presentation = pptx.Presentation(file.stream)
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                file_content += shape.text + "\n"
                except Exception as e:
                    print(f"Error reading PowerPoint file: {e}")
                    return jsonify({"reply": f"Sorry, I couldn't read the PowerPoint file '{filename}'."}), 400
            
            file_prompt_prefix = f"Based on the content of the file '{filename}' below, please answer my question.\n\nFile Content:\n---\n{file_content}\n---\n\nQuestion: "
        
        elif file_ext in media_based_ext:
            # For media files, just mention the file name as we can't process content yet
            file_prompt_prefix = f"I have uploaded a file named '{filename}'. "

        else: # Handle other/unknown file types
            file_prompt_prefix = f"I have uploaded a file named '{filename}'. "

        # Prepend file information to the user's message
        user_message = f"{file_prompt_prefix}{user_message}"

    # --- Premium Model Check ---
    if requested_model.startswith('gemini/') and not current_user.is_premium:
        error_message = "This model is only available for Premium users. Please upgrade your plan to use it."
        return jsonify({"reply": error_message, "error": True}), 403 # 403 Forbidden
    
    # --- OpenRouter Premium Model Check ---
    if requested_model.startswith('openrouter/') and not current_user.is_premium:
        error_message = "This model is only available for Premium users. Please upgrade your plan to use it."
        return jsonify({"reply": error_message, "error": True}), 403 # 403 Forbidden


    # --- Rate Limiting Check ---
    now = datetime.datetime.now(UTC)
    user_id = ObjectId(current_user.id)

    if current_user.is_premium:
        limit = 10
        window = timedelta(hours=2)
        window_start = now - window
        plan_name = "Premium"
        window_str = "2 hours"
    else:
        limit = 3
        window = timedelta(hours=3)
        window_start = now - window
        plan_name = "Free"
        window_str = "3 hours"

    prompt_count = mongo.db.prompt_logs.count_documents({
        "user_id": user_id,
        "timestamp": {"$gte": window_start}
    })

    if prompt_count >= limit:
        error_message = f"You have reached your limit of {limit} prompts per {window_str} on the {plan_name} plan. Please try again later or upgrade your plan."
        return jsonify({"reply": error_message, "error": True}), 429 # 429 Too Many Requests

    # --- Token Limit Check ---
    # A simple approximation: 1 word ≈ 1 token. For higher accuracy, a library like 'tiktoken' would be used.
    prompt_tokens = len(user_message.split())
    
    if current_user.is_premium:
        token_limit = 500
    else:
        token_limit = 250

    if prompt_tokens > token_limit:
        plan_name = "Premium" if current_user.is_premium else "Free"
        error_message = f"Your prompt is too long ({prompt_tokens} tokens). As a {plan_name} user, your limit is {token_limit} tokens per prompt."
        # Using 413 Payload Too Large status code
        return jsonify({"reply": error_message, "error": True}), 413 # 413 Payload Too Large

    title_updated = False
    # 1. Get previous messages if it's an existing chat
    history = []
    if chat_id:
        chat = mongo.db.chats.find_one({"_id": ObjectId(chat_id), "user_id": ObjectId(current_user.id)})
        if chat:
            history = chat.get("messages", [])

    # Log the prompt before calling the API
    mongo.db.prompt_logs.insert_one({"user_id": user_id, "timestamp": now})

    bot_response = ""
    model_used_by_api = requested_model # Default to requested model

    try:
        # --- API Call Logic: Gemini or Groq ---
        if requested_model.startswith('gemini/'):
            raw_model_name = requested_model.split('/')[-1]
            
            # Map legacy or unavailable model names to available ones in the workspace
            model_mappings = {
                'gemini-1.5-pro-latest': 'gemini-pro-latest',
                'gemini-1.5-flash-latest': 'gemini-flash-latest',
                'gemini-1.5-pro': 'gemini-pro-latest',
                'gemini-1.5-flash': 'gemini-flash-latest'
            }
            gemini_model_name = model_mappings.get(raw_model_name, raw_model_name)
            client = get_gemini_client()
            
            # Convert our history to the format google-genai expects
            gemini_history = []
            for msg in history:
                role = 'user' if msg['role'] == 'user' else 'model'
                gemini_history.append(
                    types.Content(
                        role=role,
                        parts=[types.Part.from_text(text=msg['content'])]
                    )
                )
            
            print(f"Sending request to Gemini via google-genai: '{gemini_model_name}' (mapped from '{raw_model_name}')")
            chat_session = client.chats.create(model=gemini_model_name, history=gemini_history)
            response = chat_session.send_message(user_message)
            bot_response = response.text
            model_used_by_api = f"gemini/{gemini_model_name}"
            print(f"Received response from Gemini. Model used: '{model_used_by_api}'")

        elif requested_model.startswith('openrouter/'):
            try:
                from openai import OpenAI
            except ImportError:
                print("ERROR: 'openai' library not installed. Please run 'pip install openai'")
                return jsonify({"reply": "Server configuration error: The 'openai' library is missing."}), 500

            openrouter_key = os.environ.get("OPENROUTER_API_KEY")
            if not openrouter_key:
                 return jsonify({"reply": "Server configuration error: OpenRouter API key is not set."}), 500

            client = OpenAI(
                base_url="https://openrouter.ai/api/v1",
                api_key=openrouter_key,
            )

            # Extract model name, removing prefix and any tags like ':free'
            openrouter_model_name = requested_model.split('/')[-1].split(':')[0]
            print(f"Sending request to OpenRouter API with model: '{openrouter_model_name}'")

            completion = client.chat.completions.create(
                model=openrouter_model_name,
                messages=[
                    {
                        "role": "system",
                        "content": "You are a helpful AI assistant."
                    },
                    *history,
                    {"role": "user", "content": user_message},
                ],
                # site_url and user_id can be passed for analytics on OpenRouter
                extra_headers={"HTTP-Referer": request.host_url, "X-Title": "ConvAI"}
            )
            bot_response = completion.choices[0].message.content
            model_used_by_api = f"openrouter/{openrouter_model_name}"
            print(f"Received response from OpenRouter. Model used: '{model_used_by_api}'")

        else: # Default to Groq
            client = Groq(api_key=os.environ.get("GROQ_API_KEY"))
            print(f"Sending request to Groq API with model: '{model}'")
            chat_completion = client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are ConvAI, a helpful and friendly AI assistant. Your name is ConvAI. When asked who you are, you must introduce yourself as ConvAI. Do not reveal that you are based on another model like Groq or Compound. Format your responses using Markdown, including tables, lists, and code blocks when appropriate."
                    },
                    *history, # Unpack previous messages
                    {
                        "role": "user",
                        "content": user_message,
                    }
                ],
                model=model,
            )
            bot_response = chat_completion.choices[0].message.content
            model_used_by_api = chat_completion.model
            print(f"Received response from Groq. Model used: '{model_used_by_api}'")

        # 2. Save the conversation to the database
        new_messages = [
            {"role": "user", "content": user_message},
            {"role": "assistant", "content": bot_response}
        ]
        
        # For Gemini, we need to adjust the saved history to match our app's format
        if requested_model.startswith('gemini/') or requested_model.startswith('openrouter/'):
            # Gemini's history is already updated in the chat_session object.
            # We just need to save our app-formatted messages.
            pass # The new_messages list is already correct.


        if chat_id:
            # Append to existing chat
            mongo.db.chats.update_one(
                {"_id": ObjectId(chat_id)},
                {
                    "$push": {"messages": {"$each": new_messages}},
                    "$set": {"updated_at": datetime.datetime.now(UTC)}
                }
            )
            return jsonify({"reply": bot_response, "chat_id": chat_id, "model_used": model_used_by_api, "title_updated": title_updated})
        else:
            # Create a new chat
            # If a file was uploaded, use its name as the title
            if file:
                title = file.filename[:50]
            else:
                # Use first 50 chars of message as title
                title = user_message.split('\n')[0][:50]

            now = datetime.datetime.now(UTC)
            new_chat_id = mongo.db.chats.insert_one({
                "user_id": ObjectId(current_user.id),
                "title": title,
                "created_at": now,
                "updated_at": now,
                "messages": new_messages
            }).inserted_id
            return jsonify({"reply": bot_response, "chat_id": str(new_chat_id), "title": title, "model_used": model_used_by_api})

    except Exception as e:
        api_name = "Gemini" if requested_model.startswith('gemini/') else "OpenRouter" if requested_model.startswith('openrouter/') else "Groq"
        print(f"Error calling {api_name} API: {e}")
        bot_response = f"Sorry, I'm having trouble connecting to the {api_name} service right now. Please try again later."
        return jsonify({"reply": bot_response}), 500
    finally:
        print("-" * 50 + "\n")

@app.route('/api/chats', methods=['GET'])
@login_required
def get_chats():
    """Fetches all chat sessions for the current user."""
    chats = mongo.db.chats.find({"user_id": ObjectId(current_user.id)}).sort("updated_at", -1)
    chat_list = [{"id": str(chat["_id"]), "title": chat["title"]} for chat in chats]
    return jsonify(chat_list)

@app.route('/api/chat/<chat_id>', methods=['GET'])
@login_required
def get_chat_messages(chat_id):
    """Fetches all messages for a specific chat session."""
    chat = mongo.db.chats.find_one({"_id": ObjectId(chat_id), "user_id": ObjectId(current_user.id)})
    if not chat:
        return jsonify({"error": "Chat not found or access denied"}), 404
    return jsonify(chat.get("messages", []))

@app.route('/api/chat/<chat_id>', methods=['DELETE'])
@login_required
def delete_chat(chat_id):
    """Deletes a chat session."""
    result = mongo.db.chats.delete_one({"_id": ObjectId(chat_id), "user_id": ObjectId(current_user.id)})
    if result.deleted_count == 1:
        return jsonify({"success": True}), 200
    else:
        return jsonify({"error": "Chat not found or access denied"}), 404

@app.route('/api/chats/clear', methods=['POST'])
@login_required
def clear_all_chats():
    """Deletes all chat sessions and prompt logs for the current user."""
    user_id = ObjectId(current_user.id)
    
    # Delete all chats for the user
    mongo.db.chats.delete_many({"user_id": user_id})
    # Also clear their prompt logs to reset rate limiting
    mongo.db.prompt_logs.delete_many({"user_id": user_id})
    
    return jsonify({"success": True, "message": "All chats have been cleared."}), 200

@app.route('/api/chat/<chat_id>', methods=['PUT'])
@login_required
def rename_chat(chat_id):
    """Renames a chat session."""
    data = request.get_json()
    new_title = data.get('title')
    if not new_title:
        return jsonify({"error": "New title not provided"}), 400

    result = mongo.db.chats.update_one(
        {"_id": ObjectId(chat_id), "user_id": ObjectId(current_user.id)},
        {"$set": {"title": new_title}}
    )
    if result.matched_count == 1:
        return jsonify({"success": True, "title": new_title}), 200
    else:
        return jsonify({"error": "Chat not found or access denied"}), 404
    
@app.route('/api/chat/<chat_id>/share', methods=['POST'])
@login_required
def share_chat(chat_id):
    """Marks a chat as public and returns the shareable link."""
    result = mongo.db.chats.update_one(
        {"_id": ObjectId(chat_id), "user_id": ObjectId(current_user.id)},
        {"$set": {"is_public": True, "updated_at": datetime.datetime.now(UTC)}}
    )
    if result.matched_count == 1:
        share_url = url_for('shared_chat', chat_id=chat_id, _external=True)
        return jsonify({"success": True, "share_url": share_url}), 200
    else:
        return jsonify({"error": "Chat not found or access denied"}), 404

@app.route('/share/<chat_id>')
def shared_chat(chat_id):
    """Displays a publicly shared chat."""
    try:
        chat = mongo.db.chats.find_one({"_id": ObjectId(chat_id), "is_public": True})
        if not chat:
            return render_template('share.html', error="This chat is not public or does not exist."), 404
        
        user = mongo.db.users.find_one({"_id": chat["user_id"]})
        return render_template('share.html', chat=chat, owner_username=user.get('username', 'Anonymous'))
    except Exception:
        return render_template('share.html', error="Invalid chat link."), 404

@app.route('/api/chat/rate', methods=['POST'])
@login_required
def rate_message():
    """Stores user feedback for a specific message."""
    data = request.get_json()
    chat_id = data.get('chat_id')
    rating = data.get('rating') # 'good' or 'bad'

    if not all([chat_id, rating]):
        return jsonify({"error": "Missing required fields"}), 400

    mongo.db.ratings.insert_one({
        "user_id": ObjectId(current_user.id),
        "chat_id": ObjectId(chat_id),
        "rating": rating,
        "rated_at": datetime.datetime.now(UTC)
    })

    return jsonify({"success": True}), 200

# --- Scheduled Tasks ---
def prune_old_chat_messages():
    """
    Finds chats that haven't been updated in 30 days and clears their message history,
    leaving a system message.
    """
    with app.app_context():
        print("Running scheduled task: Pruning old chat messages...")
        thirty_days_ago = datetime.datetime.now(UTC) - timedelta(days=30)

        # Find chats that have not been updated in the last 30 days
        old_chats = mongo.db.chats.find({"updated_at": {"$lt": thirty_days_ago}})

        for chat in old_chats:
            print(f"Pruning chat ID: {chat['_id']}")
            # Update the chat: clear messages and add a system note, and update updated_at
            mongo.db.chats.update_one(
                {"_id": chat["_id"]},
                {"$set": {
                    "messages": [{
                        "role": "system",
                        "content": "Chat history older than 30 days has been cleared to save space."
                    }],
                    "updated_at": datetime.datetime.now(UTC)
                }}
            )

def check_subscription_expiry():
    """
    Finds users with expired premium subscriptions, downgrades them, and sends an email.
    """
    with app.app_context():
        print("Running scheduled task: Checking for expired subscriptions...")
        # start_date from MongoDB is returned timezone-naive (naive UTC).
        # We use a naive UTC datetime to compare cleanly without TypeError.
        now = datetime.datetime.now(UTC).replace(tzinfo=None)
        
        # Find all premium users
        premium_users = mongo.db.users.find({"is_premium": True})

        for user in premium_users:
            start_date = user.get("subscription_start_date")
            plan_type = user.get("plan_type")

            if not start_date or not plan_type:
                continue # Skip users with incomplete subscription data

            if plan_type == 'monthly':
                expiry_date = start_date + timedelta(days=30)
            elif plan_type == 'yearly':
                expiry_date = start_date + timedelta(days=365)
            else:
                continue # Unknown plan type

            if now > expiry_date:
                user_email = user.get("email") or user.get("username")
                print(f"Subscription expired for user {user['_id']} ({user_email}). Downgrading.")
                # Downgrade user
                mongo.db.users.update_one(
                    {"_id": user["_id"]},
                    {
                        "$set": {"is_premium": False},
                        "$unset": {"plan_type": "", "subscription_start_date": ""}
                    }
                )
                # Send expiry email
                try:
                    msg = Message('Your ConvAI Premium Subscription Has Expired', recipients=[user_email])
                    msg.html = render_template('subscription_expired.html', username=user['username'])
                    mail.send(msg)
                except Exception as e:
                    print(f"Error sending subscription expiry email to {user_email}: {e}")

def send_welcome_email(user_id):
    """Function to send a welcome email to a new user."""
    with app.app_context():
        user = mongo.db.users.find_one({"_id": ObjectId(user_id)})
        if user:
            try:
                msg = Message('Welcome to ConvAI!', recipients=[user['email']])
                msg.html = render_template('email_welcome.html', username=user['username'])
                mail.send(msg)
                print(f"Welcome email sent to {user['email']}")
            except Exception as e:
                print(f"Failed to send welcome email to {user['email']}: {e}")



scheduler = BackgroundScheduler(daemon=True)

def start_scheduler():
    if not scheduler.running:
        scheduler.add_job(
            prune_old_chat_messages,
            'interval',
            days=1,
            id='prune_old_chat_messages',
            replace_existing=True
        )
        scheduler.add_job(
            check_subscription_expiry,
            'interval',
            days=1, # Run once a day
            id='check_subscription_expiry',
            replace_existing=True
        )
        scheduler.start()

def migrate_database():
    """Ensure all existing users have both 'email' and 'username' fields, fixing old data format."""
    with app.app_context():
        users = list(mongo.db.users.find())
        print(f"Checking {len(users)} users for migration...")
        for user in users:
            update_fields = {}
            
            # If email is missing
            if 'email' not in user:
                # If username is actually an email, copy it
                if '@' in user.get('username', ''):
                    email_val = user['username'].lower().strip()
                    update_fields['email'] = email_val
                    
                    # Extract a clean username if not already taken
                    clean_username = email_val.split('@')[0]
                    if not mongo.db.users.find_one({'username': clean_username}):
                        update_fields['username'] = clean_username
                else:
                    # Fallback default email
                    update_fields['email'] = f"{user.get('username')}@example.com".lower()
            
            # Save updates
            if update_fields:
                mongo.db.users.update_one({'_id': user['_id']}, {'$set': update_fields})
                print(f"Successfully migrated user {user['_id']}: {update_fields}")

def create_database_indexes():
    """Creates database indexes to optimize queries and handle high concurrent load."""
    with app.app_context():
        print("Creating database indexes for performance optimization...")
        try:
            # Index users by username and email
            mongo.db.users.create_index("username", unique=True)
        except Exception as e:
            print(f"Note: Could not create unique index on username (probably duplicates exist): {e}")
            mongo.db.users.create_index("username")

        try:
            mongo.db.users.create_index("email", unique=True)
        except Exception as e:
            print(f"Note: Could not create unique index on email (probably duplicates exist): {e}")
            mongo.db.users.create_index("email")
        
        # Compound index for chats to optimize sidebar rendering and sorting
        mongo.db.chats.create_index([("user_id", 1), ("updated_at", -1)])
        
        # Compound index for prompt logs rate limiting
        mongo.db.prompt_logs.create_index([("user_id", 1), ("timestamp", 1)])
        
        # Compound index for payments logging
        mongo.db.payments.create_index([("user_id", 1), ("timestamp", -1)])
        
        # TTL Index to automatically delete prompt logs older than 24 hours (86400 seconds)
        mongo.db.prompt_logs.create_index("timestamp", expireAfterSeconds=86400)
        
        print("Database indexes created/verified successfully.")

# --- Main Execution ---
if __name__ == '__main__':
    migrate_database()
    create_database_indexes()
    if os.environ.get("WERKZEUG_RUN_MAIN") == "true":
        start_scheduler()

    app.run(debug=True, use_reloader=True)